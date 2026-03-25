import PyPDF2
import os
from pptx import Presentation

def read_pdf(filename):
    print(f"--- Reading PDF: {filename} ---")
    try:
        with open(filename, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            # Read the first 2 pages as a sample
            num_pages = len(reader.pages)
            print(f"Total pages: {num_pages}")
            
            for i in range(min(2, num_pages)):
                print(f"\n[Page {i+1}]")
                print(reader.pages[i].extract_text())
    except Exception as e:
        print(f"Error reading {filename}: {e}")
    print("\n" + "="*30 + "\n")

def read_pptx(filename):
    print(f"--- Reading PPTX: {filename} ---")
    try:
        prs = Presentation(filename)
        print(f"Total slides: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            if i >= 2: break # Limit to first 2 slides
            print(f"\n[Slide {i+1}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except Exception as e:
        print(f"Error reading {filename}: {e}")
    print("\n" + "="*30 + "\n")

files = [
    "RWA - FastOracle.pdf",
    "2023-ACM-CS-Connect API with Blockchain A Surveyon Blockchain OracleImplementation.pdf",
    "资产代币化流程图-ZXY-7.14.pptx"
]

for f in files:
    if os.path.exists(f):
        if f.endswith('.pdf'):
            read_pdf(f)
        elif f.endswith('.pptx'):
            read_pptx(f)
    else:
        print(f"File not found: {f}")
